from pathlib import Path
import math

import matplotlib.pyplot as plt
import numpy as np
from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches as PptInches, Pt as PptPt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "deliverables"
ASSETS = OUT / "assets"
REPORT = OUT / "Smart_Application_AI_ML_Report.docx"
DECK = OUT / "Smart_Application_AI_ML_Presentation.pptx"

TITLE = "Smart Application in Artificial Intelligence and Machine Learning: Study and Understanding"
PROFESSOR = "EL MKHALET MOUNA"
STUDENT = "To be completed"


def add_page_number(paragraph):
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = paragraph.add_run()
    fld_char1 = OxmlElement("w:fldChar")
    fld_char1.set(qn("w:fldCharType"), "begin")
    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = "PAGE"
    fld_char2 = OxmlElement("w:fldChar")
    fld_char2.set(qn("w:fldCharType"), "end")
    run._r.append(fld_char1)
    run._r.append(instr)
    run._r.append(fld_char2)


def set_cell_text(cell, text, bold=False):
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
    cell.text = ""
    p = cell.paragraphs[0]
    p.paragraph_format.space_after = Pt(0)
    r = p.add_run(text)
    r.bold = bold
    r.font.name = "Times New Roman"
    r.font.size = Pt(10)


def shade_cell(cell, fill):
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def add_table(doc, headers, rows, widths=None):
    table = doc.add_table(rows=1, cols=len(headers))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    for i, h in enumerate(headers):
        set_cell_text(table.rows[0].cells[i], h, True)
        shade_cell(table.rows[0].cells[i], "D9EAF7")
    for row in rows:
        cells = table.add_row().cells
        for i, value in enumerate(row):
            set_cell_text(cells[i], str(value))
    if widths:
        for row in table.rows:
            for idx, width in enumerate(widths):
                row.cells[idx].width = Cm(width)
    doc.add_paragraph()
    return table


def add_caption(doc, text):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(text)
    r.italic = True
    r.font.size = Pt(10)


def add_body_paragraph(doc, text):
    p = doc.add_paragraph(text)
    p.paragraph_format.line_spacing = 1.5
    p.paragraph_format.space_after = Pt(6)
    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    return p


def add_bullets(doc, items):
    for item in items:
        p = doc.add_paragraph(style="List Bullet")
        p.paragraph_format.line_spacing = 1.5
        p.add_run(item)


def create_assets():
    ASSETS.mkdir(parents=True, exist_ok=True)
    blue = "#2563eb"
    green = "#16a34a"
    orange = "#f97316"
    ink = "#0f172a"

    fig, ax = plt.subplots(figsize=(8, 3.7), dpi=180)
    ax.axis("off")
    steps = ["Random data", "Preprocessing", "ML model", "Metrics", "Visualization"]
    xs = np.linspace(0.08, 0.92, len(steps))
    for i, (x, label) in enumerate(zip(xs, steps)):
        ax.add_patch(plt.Rectangle((x - 0.08, 0.38), 0.16, 0.24, color="#e0f2fe", ec=blue, lw=1.4))
        ax.text(x, 0.5, label, ha="center", va="center", fontsize=9, color=ink, weight="bold")
        if i < len(steps) - 1:
            ax.annotate("", xy=(xs[i + 1] - 0.09, 0.5), xytext=(x + 0.09, 0.5),
                        arrowprops=dict(arrowstyle="->", color=ink, lw=1.4))
    ax.set_title("Machine Learning workflow implemented in the application", fontsize=12, weight="bold")
    fig.tight_layout()
    fig.savefig(ASSETS / "ml_workflow.png", bbox_inches="tight")
    plt.close(fig)

    fig, ax = plt.subplots(figsize=(8, 4.1), dpi=180)
    ax.axis("off")
    boxes = {
        "Main desktop shell\nCustomTkinter": (0.36, 0.72, 0.28, 0.16, blue),
        "Six ML tabs\nRegression, K-Means,\nRF, ARIMA, MLP, CV": (0.06, 0.38, 0.28, 0.2, green),
        "Data generator\nnumpy.random.uniform": (0.36, 0.38, 0.28, 0.2, orange),
        "Visualization layer\nMatplotlib + TkAgg": (0.66, 0.38, 0.28, 0.2, "#a855f7"),
        "Desktop executable\nPyInstaller": (0.36, 0.08, 0.28, 0.16, "#64748b"),
    }
    for text, (x, y, w, h, color) in boxes.items():
        ax.add_patch(plt.Rectangle((x, y), w, h, color="#f8fafc", ec=color, lw=2))
        ax.text(x + w / 2, y + h / 2, text, ha="center", va="center", fontsize=9, color=ink)
    arrows = [((0.5, 0.72), (0.2, 0.58)), ((0.5, 0.72), (0.5, 0.58)),
              ((0.5, 0.72), (0.8, 0.58)), ((0.5, 0.38), (0.5, 0.24))]
    for start, end in arrows:
        ax.annotate("", xy=end, xytext=start, arrowprops=dict(arrowstyle="->", color=ink, lw=1.2))
    ax.set_title("Desktop application architecture", fontsize=12, weight="bold")
    fig.tight_layout()
    fig.savefig(ASSETS / "app_architecture.png", bbox_inches="tight")
    plt.close(fig)

    methods = ["SVM", "Naive\nBayes", "GBM", "PCA", "DBSCAN", "Apriori"]
    scores = [4, 3, 5, 3, 4, 3]
    fig, ax = plt.subplots(figsize=(7.5, 3.8), dpi=180)
    bars = ax.bar(methods, scores, color=[blue, green, orange, "#a855f7", "#14b8a6", "#f59e0b"])
    ax.set_ylim(0, 5.5)
    ax.set_ylabel("Pedagogical relevance")
    ax.set_title("Additional machine learning methods discussed in the report", weight="bold")
    ax.grid(axis="y", alpha=0.25)
    for bar, score in zip(bars, scores):
        ax.text(bar.get_x() + bar.get_width() / 2, score + 0.08, str(score), ha="center", fontsize=9)
    fig.tight_layout()
    fig.savefig(ASSETS / "methods_overview.png", bbox_inches="tight")
    plt.close(fig)


def configure_doc(doc):
    section = doc.sections[0]
    section.top_margin = Cm(2.5)
    section.bottom_margin = Cm(2.5)
    section.left_margin = Cm(2.5)
    section.right_margin = Cm(2.5)

    styles = doc.styles
    normal = styles["Normal"]
    normal.font.name = "Times New Roman"
    normal.font.size = Pt(12)
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "Times New Roman")
    normal.paragraph_format.line_spacing = 1.5
    normal.paragraph_format.space_after = Pt(6)

    for name, size, color in [
        ("Title", 18, "0F172A"),
        ("Heading 1", 14, "1D4ED8"),
        ("Heading 2", 13, "0F172A"),
        ("Heading 3", 12, "334155"),
    ]:
        style = styles[name]
        style.font.name = "Times New Roman"
        style._element.rPr.rFonts.set(qn("w:eastAsia"), "Times New Roman")
        style.font.size = Pt(size)
        style.font.color.rgb = RGBColor.from_string(color)
        style.font.bold = True
        style.paragraph_format.space_before = Pt(12)
        style.paragraph_format.space_after = Pt(6)
        style.paragraph_format.line_spacing = 1.5


INTRO_PARAS = [
    "Artificial intelligence is the scientific and technical field that designs systems able to perceive, reason, learn, decide and act with a level of autonomy. Its objective is not simply to imitate human intelligence, but to create computational processes capable of extracting value from data, supporting decisions, automating repetitive tasks and improving complex workflows. In modern organizations, AI is therefore both a technological discipline and a strategic capability.",
    "Machine learning is one of the central components of artificial intelligence. Instead of explicitly programming every rule, a machine learning system learns patterns from examples. The quality of the data, the selected algorithm, the evaluation procedure and the interpretation of results determine whether the system is useful. This project studies that logic through a desktop application that generates random data, trains several models, displays metrics and visualizes the behavior of each method.",
    "The interest of artificial intelligence comes from its ability to process information at a scale that is difficult for humans to manage manually. It helps detect fraud, classify documents, forecast demand, recommend products, monitor networks, identify anomalies, assist medical diagnosis and optimize industrial operations. Its value is particularly strong when data volume, speed or complexity exceeds classical manual analysis.",
    "The project also has a pedagogical objective. By using randomly generated data, the application isolates the behavior of algorithms from the constraints of external datasets. The user can change ranges, noise, number of clusters, model parameters or validation folds and directly observe the effect on graphs and metrics. This interactive approach makes the theory of machine learning more concrete.",
    "Artificial intelligence is studied because it is becoming a transversal competence for development, cybersecurity and intelligent systems. Developers need AI to create smarter software, cybersecurity analysts need it to detect threats and anomalies, and AI specialists need it to design models that are accurate, explainable and responsible. The following report presents the foundations, additional methods, application domains, developed executable application, limits and perspectives.",
]


ML_METHODS = [
    ("Support Vector Machines", "SVM constructs a decision boundary that maximizes the margin between classes. With kernels, it can model non-linear separations by projecting data into a higher-dimensional feature space. SVM is effective for medium-sized datasets and high-dimensional representations such as text features or cybersecurity indicators."),
    ("Naive Bayes", "Naive Bayes is a probabilistic classifier based on Bayes theorem and a simplifying independence assumption between features. Even if this assumption is rarely perfectly true, the method is fast, robust and useful for spam filtering, sentiment analysis and document classification."),
    ("Gradient Boosting Machines", "Gradient boosting builds an ensemble of weak learners, often decision trees, where each new model focuses on the errors of the previous models. Algorithms such as XGBoost, LightGBM and CatBoost are powerful for structured data and often achieve excellent predictive performance."),
    ("Principal Component Analysis", "PCA is a dimensionality reduction method that transforms correlated variables into principal components. It is not a classifier by itself, but it helps visualize high-dimensional data, reduce noise and prepare features before applying supervised or unsupervised models."),
    ("DBSCAN", "DBSCAN is a density-based clustering algorithm. Unlike K-Means, it does not require the number of clusters in advance and can identify outliers as noise. It is useful when clusters have irregular shapes or when anomaly detection is important."),
    ("Apriori and Association Rules", "Apriori discovers frequent itemsets and association rules. It is widely known in market basket analysis but can also reveal relationships between events, alerts, user actions or system states."),
    ("Hidden Markov Models", "Hidden Markov Models represent systems where the observed sequence depends on hidden states. They are useful for temporal classification, speech recognition, activity modeling and event sequence analysis."),
    ("Autoencoders", "Autoencoders are neural networks trained to reconstruct their input through a compressed internal representation. They can be used for dimensionality reduction, denoising and anomaly detection when reconstruction error becomes a signal."),
]


APP_METHODS = [
    ("Linear Regression", "The application presents a multiple linear regression in 3D. Random variables X1 and X2 are generated, a target Y is built with noise, and the model estimates coefficients and intercept. The user reads R2 and RMSE to evaluate quality."),
    ("K-Means Clustering 3D", "The clustering tab generates three random variables, applies K-Means and displays clusters in a 3D scatter plot. Centers are shown explicitly and the silhouette score provides an internal measure of cohesion and separation."),
    ("Random Forest", "The Random Forest tab generates a classification dataset, trains an ensemble of decision trees and displays accuracy and feature importances. This illustrates ensemble learning and interpretability through variable contribution."),
    ("ARIMA Time Series", "The time-series tab creates a bounded random walk, trains an ARIMA model on the training segment and forecasts future values. It displays the training curve, real test values, forecast, RMSE and AIC."),
    ("Neural Network", "The neural-network tab uses MLPRegressor, standardizes data, trains a multilayer perceptron and displays both the fitted curve and the loss curve. This helps connect model architecture to optimization behavior."),
    ("Cross Validation", "The validation tab compares Random Forest, neural network, SVM and decision tree classifiers through cross_val_score. It shows mean accuracy, standard deviation and performance evolution across folds."),
]


def add_method_section(doc):
    doc.add_heading("Chapter 1 - Presentation of Machine Learning Methods", level=1)
    add_body_paragraph(doc, "Machine learning methods can be grouped into supervised learning, unsupervised learning, semi-supervised learning, reinforcement learning and representation learning. Their efficiency depends on the nature of the problem, the structure of the data, the available computational resources and the evaluation protocol. A method is not efficient in isolation; it becomes efficient when its assumptions match the real problem.")
    add_table(doc, ["Family", "Goal", "Examples", "Typical metric"], [
        ["Supervised learning", "Predict a label or value from examples", "SVM, Gradient Boosting, Naive Bayes", "Accuracy, F1, RMSE"],
        ["Unsupervised learning", "Discover structure without labels", "PCA, DBSCAN, K-Means", "Silhouette, explained variance"],
        ["Temporal modeling", "Model sequences and time dependence", "HMM, ARIMA, LSTM", "AIC, RMSE, likelihood"],
        ["Representation learning", "Learn compact internal features", "Autoencoders, embeddings", "Reconstruction error"],
    ], [3.2, 4.4, 4.8, 3.5])
    doc.add_picture(str(ASSETS / "methods_overview.png"), width=Inches(5.9))
    add_caption(doc, "Figure 1. Additional methods discussed beyond the implemented application modules.")
    doc.add_heading("Scientific comparison criteria", level=2)
    for title, text in [
        ("Accuracy and generalization", "A model is not scientifically valid because it gives a good result once. It must generalize to new examples and remain stable when the sample changes. This is why validation strategies such as train/test split and cross-validation are essential. They transform a simple demonstration into a controlled experiment."),
        ("Complexity and interpretability", "A more complex algorithm is not automatically better. Neural networks and boosting models can capture sophisticated relationships, but their internal behavior is harder to explain. Linear models and decision trees may be less powerful in some cases, but they provide clearer explanations. In academic work, the best model is the one that balances performance and understanding."),
        ("Data assumptions", "Each method carries assumptions. K-Means assumes that clusters can be represented by centers and distances. Naive Bayes assumes conditional independence. Linear regression assumes a linear relation and residual behavior. If assumptions are ignored, the model can produce visually convincing but scientifically weak results."),
        ("Computational cost", "Efficiency also includes execution time and resource consumption. A method that is accurate but too slow may be unsuitable for real-time cybersecurity monitoring. A lightweight model may be preferable for a desktop application or embedded tool. This is why the project uses scikit-learn and statsmodels, which are appropriate for educational and desktop use."),
        ("Robustness", "A robust method tolerates noise, missing values, outliers and moderate changes in data distribution. Robustness is important because real data is rarely clean. The project uses random noise in several tabs to show that model behavior changes when the data becomes less regular."),
        ("Communication value", "A model must be communicated to a user. Graphs, metrics, captions and concise explanations help transform technical results into understandable knowledge. The application uses visual output in every module, which supports the oral defense and makes the demonstration more convincing."),
    ]:
        doc.add_heading(title, level=3)
        add_body_paragraph(doc, text)

    for name, desc in ML_METHODS:
        doc.add_heading(name, level=2)
        add_body_paragraph(doc, desc)
        add_body_paragraph(doc, f"In practice, {name} is selected when its mathematical assumptions match the structure of the available data. The method must be evaluated on unseen data because a model that performs well only on training examples does not necessarily generalize. Proper preprocessing, parameter selection and validation are therefore essential.")
        add_body_paragraph(doc, f"The pedagogical value of {name} is that it exposes an important machine learning idea: the relationship between data representation, model capacity and evaluation. When students compare this method with the algorithms implemented in the desktop application, they can understand why there is no universal algorithm and why model choice is a scientific decision.")
        doc.add_heading(f"Algorithmic reading of {name}", level=3)
        add_body_paragraph(doc, f"The algorithmic reading of {name} begins with the representation of the input variables. A raw dataset is never only a collection of numbers; it is a structured description of a problem. Features must be scaled, encoded, filtered or transformed depending on the algorithm. This step is important because many methods are sensitive to distance, variance or distributional assumptions.")
        add_body_paragraph(doc, f"After the representation step, {name} searches for a mathematical structure that minimizes error or maximizes a criterion. For supervised learning, this criterion is related to prediction quality. For unsupervised learning, it is related to structure discovery. For representation learning, it is related to the quality of a compressed or transformed view of the data.")
        add_body_paragraph(doc, f"Evaluation must be separated from training. A scientific project should avoid judging {name} only on the data used to fit it. Cross-validation, hold-out testing and repeated experiments reduce the risk of accidental conclusions. They also help compare the method with simpler baselines, which is necessary before claiming that a complex model is useful.")
        add_body_paragraph(doc, f"The advantages of {name} are meaningful only when they are connected to a domain objective. A fast method is valuable when decisions must be produced in real time. A transparent method is valuable when users require explanation. A highly accurate method is valuable when the cost of error is high. Therefore, efficiency is a combination of accuracy, robustness, interpretability and operational feasibility.")
        add_body_paragraph(doc, f"The limits of {name} must also be discussed. Some methods are sensitive to noise, some require many examples, some are hard to interpret and others need careful parameter tuning. In a professional setting, the model must be monitored after deployment because data distribution can change. This is why machine learning is not a one-time calculation but a continuous engineering process.")


def add_applications_section(doc):
    doc.add_heading("Chapter 2 - Applications of the Methods", level=1)
    add_body_paragraph(doc, "Machine learning is transversal. The same mathematical families can support software development, cybersecurity and artificial intelligence, but the objective and the constraints change from one domain to another. Development focuses on productivity and user experience, cybersecurity focuses on detection and resilience, while artificial intelligence focuses on building intelligent behavior.")
    add_table(doc, ["Domain", "Typical data", "ML uses", "Expected value"], [
        ["Development", "Logs, code metrics, user behavior, tests", "Bug prediction, recommendation, prioritization", "Better software quality and faster delivery"],
        ["Cybersecurity", "Network flows, alerts, authentication events", "Anomaly detection, malware classification, fraud detection", "Earlier detection and lower incident impact"],
        ["Artificial Intelligence", "Images, text, signals, structured data", "Prediction, perception, planning, decision support", "Automation and intelligent services"],
    ], [3.4, 4.7, 4.9, 4.2])

    for domain, paragraphs in {
        "Development": [
            "In development, machine learning can assist the entire software lifecycle. Historical repositories contain commits, bug reports, code smells, test results and deployment logs. Supervised models can estimate which modules are risky, which tests should run first or which tickets are likely to require more effort.",
            "Recommendation models can improve developer experience by suggesting code snippets, documentation pages or similar issues. Natural language processing can summarize tickets and classify support requests. In DevOps, anomaly detection can monitor system logs and detect abnormal latency, memory consumption or failure patterns after deployment.",
            "The desktop application created in this project illustrates the same logic at a smaller scale. It shows how generated data becomes input features, how models are trained and how metrics help decide whether the output is trustworthy. This is the foundation of data-driven software engineering."
        ],
        "Cybersecurity": [
            "Cybersecurity generates large volumes of heterogeneous data: firewall events, DNS queries, authentication attempts, endpoint telemetry and vulnerability reports. Manual inspection is not enough because attackers move quickly and hide malicious behavior inside normal activity.",
            "Machine learning supports intrusion detection through anomaly detection and classification. DBSCAN and autoencoders can identify unusual behavior, Random Forest and Gradient Boosting can classify malicious traffic, while Naive Bayes and SVM can support spam and phishing detection. Time-series methods also help detect abnormal bursts or repeated attack patterns.",
            "However, cybersecurity requires caution. False positives create alert fatigue and false negatives can hide attacks. Models must be validated continuously, monitored for drift and combined with expert rules, threat intelligence and human analysis."
        ],
        "Artificial Intelligence": [
            "In artificial intelligence, machine learning is used to create systems that perceive, predict and adapt. Regression models estimate continuous values, classifiers recognize categories, clustering discovers unknown groups, neural networks learn complex representations and validation methods control reliability.",
            "The application developed in this project is a simplified laboratory of these ideas. Each tab isolates a core AI mechanism: fitting, grouping, classification, forecasting, neural optimization and model comparison. The user can manipulate data ranges and observe how the model reacts.",
            "In real AI systems, these techniques are combined with data engineering, model governance, explainability and deployment pipelines. A model becomes useful only when it is integrated into a complete workflow that delivers a decision, an action or an improved user experience."
        ],
    }.items():
        doc.add_heading(domain, level=2)
        for paragraph in paragraphs:
            add_body_paragraph(doc, paragraph)
        doc.add_heading(f"Detailed scenario in {domain}", level=3)
        add_body_paragraph(doc, f"A realistic {domain.lower()} scenario starts with data collection. The data may come from logs, user actions, sensors, network events, application telemetry or historical records. Before applying a model, the team must define the target, clean the inputs, remove inconsistent values and decide which features truly represent the problem.")
        add_body_paragraph(doc, f"The second step is model selection. Simple models are often used first because they create a baseline. More advanced models are introduced only when they provide a measurable improvement. This approach is scientific: it prevents the team from choosing a sophisticated algorithm only because it looks impressive.")
        add_body_paragraph(doc, f"The third step is interpretation. In {domain.lower()}, a prediction is useful only if the user can act on it. A developer needs to know which module is risky, a cybersecurity analyst needs to know why an event is suspicious and an AI engineer needs to know how model behavior affects the final intelligent service.")
        add_body_paragraph(doc, f"The final step is communication. Results should be presented through clear graphs, concise metrics and limitations. The desktop application follows this principle by showing each model together with its visual output and numerical performance indicators.")


def add_application_chapter(doc):
    doc.add_heading("Chapter 3 - Presentation of the Executable Application", level=1)
    add_body_paragraph(doc, "The developed project is a Windows desktop application built with Python and CustomTkinter. It is delivered as an executable file generated with PyInstaller. The application validates the required machine learning components through six independent tabs and uses random data generated with numpy.random.uniform. This choice respects the project constraint: no external CSV file is required.")
    doc.add_picture(str(ASSETS / "app_architecture.png"), width=Inches(6.1))
    add_caption(doc, "Figure 2. Architecture of the executable desktop application.")
    doc.add_picture(str(ASSETS / "ml_workflow.png"), width=Inches(6.1))
    add_caption(doc, "Figure 3. General data and model workflow.")

    add_table(doc, ["Tab", "Algorithm", "Main library", "Main output"], [
        ["Regression", "Multiple Linear Regression", "sklearn.linear_model", "3D plane, R2, RMSE"],
        ["Clustering 3D", "K-Means", "sklearn.cluster", "3D clusters, centers, silhouette"],
        ["Random Forest", "RandomForestClassifier", "sklearn.ensemble", "Accuracy, feature importance"],
        ["Time Series", "ARIMA", "statsmodels", "Forecast, RMSE, AIC"],
        ["Neural Network", "MLPRegressor", "sklearn.neural_network", "Prediction curve, loss curve"],
        ["Cross Validation", "cross_val_score", "sklearn.model_selection", "Mean accuracy and standard deviation"],
    ], [3, 4.2, 4.3, 5.1])
    doc.add_heading("Experimental protocol of the application", level=2)
    for text in [
        "The application follows a simple experimental protocol that can be explained during the defense. First, the user selects the parameter values with sliders. Second, the application generates random data using controlled ranges. Third, the selected model is trained. Fourth, the application calculates one or more metrics. Fifth, the graph is refreshed so the user can observe the result.",
        "This protocol is important because it makes the project reproducible at the level of logic. The exact random values change from one run to another, but the procedure remains the same. This is a useful way to understand stochastic experiments: the student can repeat a test and observe whether the model behavior remains coherent.",
        "The use of random generation also avoids dependency on external files. In many academic projects, missing CSV paths or incompatible encodings can prevent a demonstration from working. Here, the executable application contains the full learning experience without requiring a dataset download.",
        "The metrics were selected to match the model family. Regression uses R2 and RMSE, classification uses accuracy, clustering uses silhouette score and time-series modeling uses RMSE and AIC. This correspondence shows that evaluation is not generic; it must match the task type.",
        "The interface was improved with a consistent visual design, dark theme, readable panels and embedded charts. This matters because a final project is evaluated not only by code but also by clarity, professionalism and the ability to communicate results quickly.",
    ]:
        add_body_paragraph(doc, text)

    for name, desc in APP_METHODS:
        doc.add_heading(name, level=2)
        add_body_paragraph(doc, desc)
        add_body_paragraph(doc, "The interface separates parameters from visualization. The left panel contains sliders and execution buttons, while the right panel contains Matplotlib graphs embedded in the desktop window. This design supports experimentation: the student changes a parameter, runs the model and immediately compares visual and numerical results.")
        add_body_paragraph(doc, "The executable packaging is important because the validation document explicitly states that the application must be provided as an executable file. A project that is not executable can be sanctioned. For this reason, the final application is generated in a desktop distribution folder and can be launched by double-clicking the executable or a Windows shortcut.")
        add_body_paragraph(doc, f"From a user-experience point of view, the {name} module is designed to reduce cognitive effort. The user does not need to write code or prepare a dataset. Parameters are adjusted through sliders, the model is executed with one button and the result appears immediately. This is important for an academic demonstration because the audience can focus on the algorithmic idea rather than installation or command-line details.")
        add_body_paragraph(doc, f"From a technical point of view, the {name} module follows a reproducible structure: collect interface parameters, generate random data, fit the model, compute metrics, clear the previous graph and draw the new result. This repeated structure makes the application maintainable because each tab is independent while still sharing the same design logic.")
        add_body_paragraph(doc, f"The interpretation of the {name} result must combine graph and metric. A graph can reveal shape, dispersion, outliers or instability, while a metric summarizes performance numerically. The student should not present only numbers or only visuals; the strength of the application is precisely that it connects both forms of evidence.")


def add_limits_section(doc):
    doc.add_heading("Chapter 4 - Limits and Perspectives", level=1)
    add_body_paragraph(doc, "Machine learning has strong potential, but it also has limits. These limits are technical, methodological, ethical and organizational. A responsible project must recognize them instead of presenting AI as an automatic solution to every problem.")
    add_table(doc, ["Limit", "Impact", "Mitigation"], [
        ["Data quality", "Wrong or biased data produces unreliable predictions", "Data cleaning, validation, governance"],
        ["Overfitting", "Excellent training results but weak generalization", "Cross-validation, regularization, test sets"],
        ["Interpretability", "Complex models can be difficult to explain", "Feature importance, SHAP, simpler baselines"],
        ["Drift", "Model performance changes when real data changes", "Monitoring and retraining"],
        ["Security", "Models can be attacked or manipulated", "Robustness tests, adversarial awareness"],
        ["Ethics", "Unfair decisions can harm users", "Bias audits and transparent use"],
    ], [3.2, 5.3, 5.3])
    doc.add_heading("Critical reading of machine learning results", level=2)
    for text in [
        "A common mistake is to interpret a metric as an absolute truth. In reality, a metric is a summary produced under specific conditions. If the data changes, if the split changes or if the parameter values change, the metric can change too. Scientific interpretation therefore requires repetition and comparison.",
        "Another mistake is to confuse correlation with causality. A model may find a statistical relationship without proving that one variable causes another. This distinction is essential in AI projects because automated decisions can influence people, organizations and security processes.",
        "Visualization is useful because it reveals aspects that metrics can hide. A regression line may have a good average error while still failing in a specific region. A clustering score may look acceptable while some clusters overlap visually. A forecast may have a low RMSE but miss the direction of future movement.",
        "For this reason, the report and presentation should emphasize both quantitative and qualitative interpretation. During the defense, the student can show a graph, read the metric and explain what the result means in simple terms. This is stronger than presenting code alone.",
        "The best perspective for this project is to treat it as a foundation. It is already executable and pedagogical, and it can evolve toward a more professional tool by adding real datasets, model export, automatic report generation and explainability features.",
    ]:
        add_body_paragraph(doc, text)
    for title, text in [
        ("Data dependency", "ML systems depend heavily on representative data. In this project the data is randomly generated, which is excellent for learning algorithms, but it does not reproduce all the complexity of real business data. A future version could include controlled CSV import while preserving the random generation mode."),
        ("Model evaluation", "Metrics are essential but they can be misleading if used alone. Accuracy can hide class imbalance, RMSE can hide local errors and silhouette score does not guarantee business relevance. Evaluation must combine quantitative metrics, visualization and domain understanding."),
        ("Explainability and trust", "Users must understand why a model makes a decision. Random Forest importance and regression coefficients are easier to explain than deep neural networks. Future work could add explainability panels, confusion matrices and model comparison summaries."),
        ("Deployment perspective", "The current executable proves desktop deployment. Future improvements could add automatic graph export, report generation, persistent projects, CSV import, model saving and a bilingual interface. Another perspective is to publish the source code and executable on GitHub and present the project on LinkedIn as required by the validation document."),
    ]:
        doc.add_heading(title, level=2)
        add_body_paragraph(doc, text)
        add_body_paragraph(doc, "This perspective connects the academic objective with professional practice. A strong AI project is not only a model; it is a complete system with usability, documentation, validation, deployment and communication.")
        add_body_paragraph(doc, "Another important point is maintainability. A model that works during a classroom demonstration may need additional controls before being used in production. Error handling, version control, logging, data validation, security checks and user documentation transform a prototype into a professional tool.")
        add_body_paragraph(doc, "Finally, the human factor remains central. Machine learning supports decision-making but does not replace responsibility. The user must understand the limits of the model, verify unexpected results and communicate uncertainty. This is especially important in cybersecurity and AI systems where a wrong decision can have operational or ethical consequences.")

    doc.add_heading("Project validation and defense strategy", level=2)
    for text in [
        "The validation document insists on an executable application, a report, a concise PowerPoint presentation and a five-minute English video. These elements should not be treated separately. They are four views of the same project: the application proves implementation, the report proves understanding, the presentation proves synthesis and the video proves communication.",
        "During the defense, the first objective is to show that the executable works. The student should open the application from the Windows shortcut or executable file, present the first interface and mention the professor name, the student name and the EMSI identity requirement. This immediately answers the strict validation condition.",
        "The second objective is to explain the logic of the tabs. The student should not spend too much time on installation or code details. A stronger strategy is to open two or three representative tabs, change parameters, run the model and explain the observed output. Regression, clustering and cross-validation are good choices because they show different families of machine learning.",
        "The third objective is to connect the demonstration to the report. For example, when showing K-Means, the student can recall the difference between supervised and unsupervised learning. When showing Random Forest, the student can mention ensemble methods and feature importance. When showing cross-validation, the student can explain why evaluation must be repeated across folds.",
        "The fourth objective is to discuss limitations honestly. A strong academic presentation does not pretend that random data is equivalent to production data. Instead, it explains that random data is used here for pedagogical control and that a future version could add CSV import, model persistence and explainability tools.",
        "The final objective is professional positioning. The validation document recommends publishing the project on LinkedIn or GitHub. This is important because the project demonstrates several visible skills: Python development, desktop UI design, machine learning, visualization, packaging, reporting and English technical communication.",
        "The PowerPoint should remain concise. It should not repeat the entire report. Its role is to guide the oral explanation with essential claims, visuals and results. The report carries the detailed academic content, while the deck carries the story of the project.",
        "The English video should be rehearsed with a timer. Five minutes is short, so each part must be controlled: thirty seconds for introduction, one minute for project objective, two minutes for application modules, one minute for domains and limits, and thirty seconds for conclusion. This structure is included in the appendix.",
        "A good defense also anticipates questions. If asked why random data was used, the answer is that the project statement required random generation and that it helps compare model behavior without external file dependency. If asked why these libraries were selected, the answer is that scikit-learn, statsmodels, Matplotlib and CustomTkinter are lightweight and appropriate for an educational desktop application.",
        "If asked about improvement, the student can propose a roadmap: CSV import, graph export, confusion matrices, model saving, richer first interface, EMSI logo integration, and a GitHub release containing the executable, presentation and video. This shows that the project is complete but still open to professional evolution.",
    ]:
        add_body_paragraph(doc, text)


def add_appendices(doc):
    doc.add_heading("Conclusion", level=1)
    add_body_paragraph(doc, "This report presented a complete study and understanding of an intelligent desktop application in artificial intelligence and machine learning. It introduced AI concepts, objectives, domains, components and reasons for studying the field. It then described several machine learning methods, their efficiency and their possible use in development, cybersecurity and artificial intelligence.")
    add_body_paragraph(doc, "The executable application demonstrates the practical part of the project through six modules: regression, 3D clustering, Random Forest, ARIMA time series, neural network regression and cross-validation. The project respects the constraint of randomly generated data and provides an executable Windows application suitable for final validation.")
    add_body_paragraph(doc, "The main lesson is that machine learning is both theoretical and practical. Algorithms must be understood, implemented, evaluated and communicated. The desktop application, report, PowerPoint presentation and English video together create a complete academic and professional portfolio.")

    doc.add_heading("References", level=1)
    refs = [
        "Pedagogical project validation document, Dr. EL MKHALET MOUNA, Project Machine Learning - Artificial Intelligence.",
        "Recherche Documentaire Scientifique, Prof. Houssam Moustansir, EMSI Casablanca, 3IIR S2 2025/2026, especially slide 90 formatting recommendations.",
        "Scikit-learn documentation: supervised learning, clustering, neural networks and model validation.",
        "Statsmodels documentation: ARIMA time-series modeling.",
        "Matplotlib documentation: embedded visualizations and plotting.",
        "Python documentation: desktop application development and packaging ecosystem.",
    ]
    for ref in refs:
        p = doc.add_paragraph(style="List Number")
        p.paragraph_format.line_spacing = 1.5
        p.add_run(ref)

    doc.add_heading("Appendix A - Five-minute English video script", level=1)
    video = [
        ("0:00 - 0:30", "Hello, my name is [student name]. In this video I present my project entitled Smart Application in Artificial Intelligence and Machine Learning: Study and Understanding. The project was supervised by Professor EL MKHALET MOUNA."),
        ("0:30 - 1:10", "The objective is to understand artificial intelligence and machine learning through a desktop application. The application generates random data, trains models and displays metrics and graphs."),
        ("1:10 - 2:20", "The application contains six modules: multiple linear regression, 3D K-Means clustering, Random Forest classification, ARIMA time-series forecasting, neural network regression and cross-validation."),
        ("2:20 - 3:30", "Each module allows the user to change parameters with sliders. The user can observe the effect on the model and compare numerical results such as R2, RMSE, accuracy, silhouette score and AIC."),
        ("3:30 - 4:20", "The project is important for development, cybersecurity and artificial intelligence. In development it helps build intelligent software. In cybersecurity it supports anomaly detection. In AI it demonstrates how models learn from data."),
        ("4:20 - 5:00", "The final result is an executable desktop application, a complete report and a PowerPoint presentation. This project helped me connect theory, implementation, evaluation and professional communication."),
    ]
    add_table(doc, ["Time", "Script"], video, [3, 13])

    doc.add_heading("Appendix B - Validation checklist", level=1)
    add_bullets(doc, [
        "First interface includes EMSI logo, professor name and student name.",
        "Application is delivered as an executable Windows file.",
        "Report title matches the required title exactly.",
        "Report mentions Professor EL MKHALET MOUNA and the student name.",
        "Report follows Introduction, Chapter 1, Chapter 2, Chapter 3, Chapter 4.",
        "PowerPoint presentation is concise and highlights essential elements.",
        "Five-minute English video is prepared and structured.",
        "Final project can be published on LinkedIn or GitHub with application, presentation and video.",
    ])


def build_report():
    doc = Document()
    configure_doc(doc)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("EMSI\n")
    r.bold = True
    r.font.size = Pt(16)
    r = p.add_run(TITLE)
    r.bold = True
    r.font.size = Pt(18)
    doc.add_paragraph()
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta.add_run(f"Professor: {PROFESSOR}\n").bold = True
    meta.add_run(f"Student name: {STUDENT}\n").bold = True
    meta.add_run("Project: Machine Learning - Artificial Intelligence\n")
    meta.add_run("Academic year: 2025/2026")
    for _ in range(10):
        doc.add_paragraph()
    note = doc.add_paragraph()
    note.alignment = WD_ALIGN_PARAGRAPH.CENTER
    note.add_run("Complete project report prepared according to the validation document and slide 90 formatting guide.")
    doc.add_page_break()

    doc.add_heading("Abstract", level=1)
    add_body_paragraph(doc, "This report studies artificial intelligence and machine learning through the design and presentation of a smart desktop application. The application generates random data, applies several models, displays visualizations and provides performance metrics. The report follows the requested academic structure and includes methodological discussion, practical applications, executable application description, limitations, perspectives and validation appendices.")

    doc.add_heading("Table of Contents", level=1)
    for item in [
        "Introduction",
        "Chapter 1 - Presentation of Machine Learning Methods",
        "Chapter 2 - Applications of the Methods",
        "Chapter 3 - Presentation of the Executable Application",
        "Chapter 4 - Limits and Perspectives",
        "Conclusion",
        "References",
        "Appendices",
    ]:
        doc.add_paragraph(item)
    doc.add_page_break()

    section = doc.add_section(WD_SECTION.NEW_PAGE)
    section.top_margin = Cm(2.5)
    section.bottom_margin = Cm(2.5)
    section.left_margin = Cm(2.5)
    section.right_margin = Cm(2.5)
    add_page_number(section.footer.paragraphs[0])

    doc.add_heading("Introduction", level=1)
    for para in INTRO_PARAS:
        add_body_paragraph(doc, para)
    add_table(doc, ["AI component", "Role in intelligent systems"], [
        ["Data", "Raw material used to learn patterns and evaluate behavior."],
        ["Algorithms", "Mathematical procedures that transform data into predictions or decisions."],
        ["Models", "Learned representations produced by training algorithms."],
        ["Evaluation", "Metrics and tests used to verify reliability."],
        ["Deployment", "Integration of models into usable software or services."],
        ["Ethics", "Rules and practices that protect fairness, privacy and accountability."],
    ], [4, 12])

    add_method_section(doc)
    add_applications_section(doc)
    add_application_chapter(doc)
    add_limits_section(doc)
    add_appendices(doc)

    doc.core_properties.title = TITLE
    doc.core_properties.subject = "Artificial Intelligence and Machine Learning project report"
    doc.core_properties.author = STUDENT
    doc.save(REPORT)


def add_slide_title(slide, kicker, title):
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0.0), PptInches(0), PptInches(0.12), PptInches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = PptRGB(37, 99, 235)
    accent.line.fill.background()
    box = slide.shapes.add_textbox(PptInches(0.58), PptInches(0.35), PptInches(2.0), PptInches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = kicker.upper()
    p.font.size = PptPt(8)
    p.font.bold = True
    p.font.color.rgb = PptRGB(56, 189, 248)
    p.font.name = "Arial"
    title_box = slide.shapes.add_textbox(PptInches(0.58), PptInches(0.67), PptInches(8.7), PptInches(1.0))
    title_box.text_frame.word_wrap = True
    title_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = PptPt(22)
    p.font.bold = True
    p.font.color.rgb = PptRGB(15, 23, 42)
    p.font.name = "Arial"


def add_footer(slide, n):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0.58), PptInches(6.95), PptInches(8.6), PptInches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = PptRGB(203, 213, 225)
    line.line.fill.background()
    box = slide.shapes.add_textbox(PptInches(9.45), PptInches(6.78), PptInches(0.45), PptInches(0.2))
    p = box.text_frame.paragraphs[0]
    p.text = f"{n:02d}"
    p.font.size = PptPt(9)
    p.font.color.rgb = PptRGB(100, 116, 139)


def add_bullet_box(slide, items, x, y, w, h, font_size=13):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = PptPt(font_size)
        p.font.name = "Arial"
        p.font.color.rgb = PptRGB(30, 41, 59)
        p.space_after = PptPt(6)


def add_metric(slide, value, label, x, y, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(y), PptInches(2.15), PptInches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = PptPt(24)
    p.font.bold = True
    p.font.color.rgb = PptRGB(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = PptPt(9)
    p2.font.color.rgb = PptRGB(226, 232, 240)
    p2.alignment = PP_ALIGN.CENTER


def add_card(slide, title, body, x, y, w, h, color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = PptRGB(255, 255, 255)
    card.line.color.rgb = PptRGB(226, 232, 240)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(x), PptInches(y), PptInches(0.08), PptInches(h))
    band.fill.solid()
    band.fill.fore_color.rgb = color
    band.line.fill.background()
    t = slide.shapes.add_textbox(PptInches(x + 0.22), PptInches(y + 0.18), PptInches(w - 0.35), PptInches(0.28))
    p = t.text_frame.paragraphs[0]
    p.text = title
    p.font.size = PptPt(12)
    p.font.bold = True
    p.font.color.rgb = PptRGB(15, 23, 42)
    b = slide.shapes.add_textbox(PptInches(x + 0.22), PptInches(y + 0.52), PptInches(w - 0.35), PptInches(h - 0.62))
    b.text_frame.word_wrap = True
    p = b.text_frame.paragraphs[0]
    p.text = body
    p.font.size = PptPt(10)
    p.font.color.rgb = PptRGB(71, 85, 105)


def add_process(slide, labels, x, y, w):
    colors = [PptRGB(37, 99, 235), PptRGB(22, 163, 74), PptRGB(249, 115, 22), PptRGB(124, 58, 237)]
    gap = 0.18
    step_w = (w - gap * (len(labels) - 1)) / len(labels)
    for i, label in enumerate(labels):
        sx = x + i * (step_w + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(sx), PptInches(y), PptInches(step_w), PptInches(0.86))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[i % len(colors)]
        shape.line.fill.background()
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = PptPt(10)
        p.font.bold = True
        p.font.color.rgb = PptRGB(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER


def add_side_panel(slide, heading, lines, color=PptRGB(15, 23, 42)):
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(6.55), PptInches(1.65), PptInches(2.85), PptInches(4.7))
    panel.fill.solid()
    panel.fill.fore_color.rgb = color
    panel.line.fill.background()
    title = slide.shapes.add_textbox(PptInches(6.85), PptInches(1.95), PptInches(2.3), PptInches(0.35))
    p = title.text_frame.paragraphs[0]
    p.text = heading
    p.font.size = PptPt(13)
    p.font.bold = True
    p.font.color.rgb = PptRGB(255, 255, 255)
    box = slide.shapes.add_textbox(PptInches(6.85), PptInches(2.45), PptInches(2.25), PptInches(3.5))
    box.text_frame.word_wrap = True
    for i, line in enumerate(lines):
        p = box.text_frame.paragraphs[0] if i == 0 else box.text_frame.add_paragraph()
        p.text = line
        p.font.size = PptPt(10.5)
        p.font.color.rgb = PptRGB(226, 232, 240)
        p.space_after = PptPt(8)


def build_deck():
    prs = Presentation()
    prs.slide_width = PptInches(10)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    slides = [
        ("PROJECT", TITLE, [
            "Professor: EL MKHALET MOUNA",
            "Student: To be completed",
            "Deliverables: executable desktop app, report, PowerPoint and English video",
        ]),
        ("CONTEXT", "The project answers a strict validation requirement.", [
            "The application must be executable as a Windows .exe.",
            "The first interface must mention EMSI, the professor and the student.",
            "The report must follow the requested academic structure.",
            "The presentation must highlight the essential logic of the project.",
        ]),
        ("OBJECTIVE", "The goal is to transform ML theory into an interactive desktop laboratory.", [
            "Generate random data using controlled min/max ranges.",
            "Train several machine learning models from the same application.",
            "Display graphs and metrics so the user can interpret the model.",
            "Connect theory, implementation, evaluation and communication.",
        ]),
        ("AI OVERVIEW", "Artificial intelligence creates systems able to learn, reason and support decisions.", [
            "AI uses data, algorithms, models, evaluation and deployment workflows.",
            "It is used in recommendation, automation, forecasting, security and decision support.",
            "Its value comes from processing large, complex or fast-changing information.",
            "The project studies AI through the practical subfield of machine learning.",
        ]),
        ("ML FOUNDATION", "Machine learning learns patterns from data instead of hard-coding every rule.", [
            "Supervised learning predicts labels or numerical values.",
            "Unsupervised learning discovers hidden structure without labels.",
            "Time-series learning models order, memory and future evolution.",
            "Validation estimates whether results can generalize beyond one run.",
        ]),
        ("WORKFLOW", "Every tab follows the same scientific workflow.", [
            "1. Choose parameters with sliders.",
            "2. Generate random data with numpy.random.uniform.",
            "3. Train the selected model.",
            "4. Display metrics and visualization for interpretation.",
        ]),
        ("ARCHITECTURE", "The desktop application is modular and easy to demonstrate.", [
            "Main shell: CustomTkinter window with six independent tabs.",
            "Model layer: scikit-learn and statsmodels algorithms.",
            "Visualization layer: Matplotlib embedded with TkAgg.",
            "Packaging layer: PyInstaller produces the final Windows executable.",
        ]),
        ("DATA", "Random data is a requirement and a pedagogical advantage.", [
            "No external CSV file is needed during the defense.",
            "The user controls ranges, noise, clusters, folds and model parameters.",
            "Random generation makes repeated experiments possible.",
            "The limitation is that random data is less realistic than production datasets.",
        ]),
        ("REGRESSION", "Regression explains how a model estimates a continuous value.", [
            "The application generates X1, X2 and a target Y with random noise.",
            "LinearRegression estimates coefficients and intercept.",
            "The 3D graph shows the point cloud and the fitted regression plane.",
            "R2 and RMSE measure explanatory power and prediction error.",
        ]),
        ("CLUSTERING", "K-Means 3D shows how unsupervised learning discovers groups.", [
            "The model receives X1, X2 and X3 without class labels.",
            "K-Means places centers and assigns each point to the nearest center.",
            "The 3D visualization makes cluster separation visible.",
            "The silhouette score evaluates cohesion and separation.",
        ]),
        ("RANDOM FOREST", "Random Forest illustrates ensemble learning and feature importance.", [
            "Several decision trees are trained on random subsets of the data.",
            "The final prediction combines the trees to reduce instability.",
            "Accuracy evaluates classification performance.",
            "Feature importance shows which variables contributed most.",
        ]),
        ("TIME SERIES", "ARIMA introduces forecasting with ordered observations.", [
            "The application creates a bounded random walk as a time series.",
            "ARIMA(p,d,q) models autoregression, differencing and moving average effects.",
            "The graph separates training data, test data and forecast.",
            "RMSE and AIC help evaluate prediction error and model quality.",
        ]),
        ("NEURAL NETWORK", "The neural network tab explains non-linear learning and optimization.", [
            "MLPRegressor learns a function using hidden layers and neurons.",
            "StandardScaler normalizes the input before training.",
            "The first graph compares real values and predictions.",
            "The loss curve shows how training error evolves over iterations.",
        ]),
        ("CROSS VALIDATION", "Cross-validation avoids judging a model from one lucky split.", [
            "The dataset is divided into several folds.",
            "Each model is trained and tested several times.",
            "The application compares Random Forest, neural network, SVM and decision tree.",
            "Mean accuracy and standard deviation show performance and stability.",
        ]),
        ("ADDITIONAL METHODS", "The report studies extra methods beyond the application tabs.", [
            "SVM: builds a margin-based decision boundary.",
            "Naive Bayes: fast probabilistic classification.",
            "Gradient Boosting: combines weak learners to reduce errors.",
            "PCA: reduces dimensions and helps visualize complex data.",
        ]),
        ("MORE METHODS", "Other methods extend the project toward real AI use cases.", [
            "DBSCAN detects density-based clusters and outliers.",
            "Apriori discovers association rules between events or items.",
            "Hidden Markov Models represent hidden states in sequences.",
            "Autoencoders learn compressed representations and detect anomalies.",
        ]),
        ("DEVELOPMENT", "In development, ML improves software quality and productivity.", [
            "Bug prediction can identify risky modules before release.",
            "Log analysis can detect abnormal application behavior.",
            "Recommendation systems can suggest code, tests or documentation.",
            "DevOps monitoring can prioritize incidents and reduce downtime.",
        ]),
        ("CYBERSECURITY", "In cybersecurity, ML helps detect threats inside massive event streams.", [
            "Anomaly detection finds unusual network or authentication behavior.",
            "Classification models detect spam, phishing or malware patterns.",
            "Time-series analysis can detect bursts or repeated attack attempts.",
            "The main risk is false positives, so human validation remains necessary.",
        ]),
        ("AI OPTION", "In artificial intelligence, these methods become intelligent services.", [
            "Regression supports estimation and forecasting.",
            "Classification supports recognition and decision assistance.",
            "Clustering supports segmentation and discovery.",
            "Neural networks support complex representations and adaptation.",
        ]),
        ("INTERPRETATION", "A good ML demonstration must explain both graph and metric.", [
            "Graphs reveal shape, dispersion, outliers and instability.",
            "Metrics summarize performance but can hide local errors.",
            "The same metric is not valid for every problem type.",
            "The presenter must connect the result to the algorithmic objective.",
        ]),
        ("LIMITS", "Machine learning is powerful, but it is not magic.", [
            "Data quality strongly controls model quality.",
            "Overfitting can produce good training results but weak generalization.",
            "Complex models can be difficult to explain.",
            "Models can drift when real-world data changes.",
        ]),
        ("PERSPECTIVES", "The next version can become a stronger professional portfolio project.", [
            "Add CSV import while keeping random generation mode.",
            "Export graphs and automatically generate experiment reports.",
            "Add confusion matrices, SHAP explanations and model saving.",
            "Publish the executable, presentation and video on GitHub or LinkedIn.",
        ]),
        ("VIDEO PLAN", "The five-minute English video should be structured and timed.", [
            "0:00-0:30 introduce the title, professor and student.",
            "0:30-1:10 explain objective and AI/ML context.",
            "1:10-3:30 demonstrate the six application modules.",
            "3:30-5:00 explain domains, limits, perspectives and conclusion.",
        ]),
        ("DEFENSE", "The oral defense should prove execution, understanding and professionalism.", [
            "Start by launching the executable from the desktop shortcut.",
            "Show two or three representative tabs live.",
            "Explain why random data was required and useful.",
            "Finish with limits, improvements and publication plan.",
        ]),
        ("CONCLUSION", "The project connects research, implementation and communication.", [
            "The report proves scientific understanding.",
            "The executable proves technical implementation.",
            "The presentation proves synthesis and communication.",
            "The video and GitHub/LinkedIn publication create a visible portfolio.",
        ]),
    ]

    for idx, (kicker, title, bullets) in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = PptRGB(248, 250, 252)
        if idx == 1:
            hero = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(0), PptInches(10), PptInches(7.5))
            hero.fill.solid()
            hero.fill.fore_color.rgb = PptRGB(15, 23, 42)
            hero.line.fill.background()
            strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(0), PptInches(0.18), PptInches(7.5))
            strip.fill.solid()
            strip.fill.fore_color.rgb = PptRGB(56, 189, 248)
            strip.line.fill.background()
            title_box = slide.shapes.add_textbox(PptInches(0.75), PptInches(0.95), PptInches(8.2), PptInches(2.0))
            title_box.text_frame.word_wrap = True
            title_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            p = title_box.text_frame.paragraphs[0]
            p.text = title
            p.font.size = PptPt(28)
            p.font.bold = True
            p.font.color.rgb = PptRGB(255, 255, 255)
            add_bullet_box(slide, bullets, 0.78, 3.45, 5.45, 1.7, 14)
            for shape in slide.shapes:
                if shape.has_text_frame and shape.left < PptInches(6.2):
                    for para in shape.text_frame.paragraphs:
                        para.font.color.rgb = PptRGB(226, 232, 240)
            add_metric(slide, "6", "ML modules", 6.75, 3.2, PptRGB(37, 99, 235))
            add_metric(slide, ".exe", "desktop delivery", 6.75, 4.35, PptRGB(22, 163, 74))
            add_footer(slide, idx)
            continue
        else:
            add_slide_title(slide, kicker, title)
            if idx in (6, 7, 8):
                img = ASSETS / ("app_architecture.png" if idx == 7 else "ml_workflow.png")
                add_bullet_box(slide, bullets, 0.75, 1.75, 4.95, 3.7)
                slide.shapes.add_picture(str(img), PptInches(5.75), PptInches(2.1), width=PptInches(3.8))
            elif idx in (9, 10, 11, 12, 13, 14):
                add_bullet_box(slide, bullets[:2], 0.75, 1.75, 4.75, 1.35)
                add_process(slide, bullets[2:] if len(bullets) > 2 else bullets, 0.75, 3.65, 5.1)
                module_notes = {
                    9: ["Output: 3D regression plane", "Metrics: R2 and RMSE", "Goal: estimate a continuous value"],
                    10: ["Output: 3D clusters", "Metric: silhouette score", "Goal: discover hidden groups"],
                    11: ["Output: class prediction", "Metric: accuracy", "Bonus: feature importance"],
                    12: ["Output: forecast curve", "Metrics: RMSE and AIC", "Goal: predict future values"],
                    13: ["Output: fitted curve", "Metric: RMSE", "Bonus: loss curve"],
                    14: ["Output: model comparison", "Metrics: mean and std", "Goal: stable evaluation"],
                }
                add_side_panel(slide, "Module proof", module_notes[idx], PptRGB(30, 41, 59))
            elif idx in (15, 16):
                labels = [b.split(":")[0] for b in bullets]
                bodies = [b.split(":", 1)[1].strip() if ":" in b else b for b in bullets]
                coords = [(0.75, 1.8), (3.25, 1.8), (0.75, 4.0), (3.25, 4.0)]
                colors = [PptRGB(37, 99, 235), PptRGB(22, 163, 74), PptRGB(249, 115, 22), PptRGB(124, 58, 237)]
                for (x, y), label, body, color in zip(coords, labels, bodies, colors):
                    add_card(slide, label, body, x, y, 2.25, 1.55, color)
                add_side_panel(slide, "Why it matters", ["Shows the ML ecosystem", "Completes the report chapter", "Prepares possible defense questions"], PptRGB(37, 99, 235))
            elif idx in (17, 18, 19):
                add_bullet_box(slide, bullets, 0.75, 1.75, 5.2, 3.7)
                domain = {17: ("Development", "Quality\nProductivity\nDevOps"), 18: ("Cybersecurity", "Detection\nMonitoring\nResilience"), 19: ("AI option", "Prediction\nPerception\nDecision")}[idx]
                add_side_panel(slide, domain[0], domain[1].split("\n"), PptRGB(22, 101, 52) if idx == 17 else PptRGB(127, 29, 29) if idx == 18 else PptRGB(67, 56, 202))
            elif idx in (20, 21, 22, 23, 24, 25):
                add_bullet_box(slide, bullets, 0.75, 1.8, 5.25, 3.65)
                right_titles = {
                    20: ("Interpretation", ["Graph + metric", "Context", "Conclusion"]),
                    21: ("Risks", ["Data quality", "Overfitting", "Drift"]),
                    22: ("Roadmap", ["CSV import", "Explainability", "Publication"]),
                    23: ("Timing", ["5 minutes", "English", "Structured"]),
                    24: ("Defense", ["Launch .exe", "Show modules", "Explain limits"]),
                    25: ("Final message", ["Research", "Implementation", "Communication"]),
                }
                add_side_panel(slide, right_titles[idx][0], right_titles[idx][1], PptRGB(15, 23, 42))
            else:
                add_bullet_box(slide, bullets, 0.75, 1.75, 5.25, 3.6)
                add_side_panel(slide, "Core idea", ["Theory", "Practice", "Evaluation"], PptRGB(37, 99, 235))
        add_footer(slide, idx)

    prs.save(DECK)


def add_image_frame(slide, image_path, x, y, w, h=None):
    frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        PptInches(x - 0.05),
        PptInches(y - 0.05),
        PptInches(w + 0.1),
        PptInches((h or w * 0.58) + 0.1),
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = PptRGB(255, 255, 255)
    frame.line.color.rgb = PptRGB(203, 213, 225)
    slide.shapes.add_picture(str(image_path), PptInches(x), PptInches(y), width=PptInches(w))


def build_deck():
    """Short defense deck: designed for a maximum 7-minute oral presentation."""
    prs = Presentation()
    prs.slide_width = PptInches(10)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    screenshots = ASSETS / "screenshots"
    slides = [
        ("PROJECT", TITLE, [
            "Professor: EL MKHALET MOUNA",
            "Student: To be completed",
            "Executable desktop app, report, presentation and English video",
        ]),
        ("TABLE OF CONTENT", "Presentation roadmap", [
            "1. Foundation: AI and machine learning basics.",
            "2. Workflow: how each module runs an experiment.",
            "3. Architecture: interface, algorithms, visualization and packaging.",
            "4. Demonstration: screenshots and module results.",
            "5. Methods and applications: algorithms, domains and limits.",
            "6. Defense plan and conclusion.",
        ]),
        ("REQUIREMENTS", "The project must prove execution, understanding and communication.", [
            "Executable Windows application required for validation.",
            "Report title and structure follow the official PDF.",
            "Presentation stays concise for a maximum 7-minute defense.",
            "English video script is prepared for the final package.",
        ]),
        ("FOUNDATION", "AI and machine learning convert data into useful decisions.", [
            "AI combines data, algorithms, models, evaluation and deployment.",
            "Machine learning learns patterns instead of hard-coding every rule.",
            "The project uses random data to make experiments controlled and repeatable.",
        ]),
        ("WORKFLOW", "Each module follows one clear scientific workflow.", [
            "Choose parameters with sliders.",
            "Generate random data using numpy.random.uniform.",
            "Train the model and calculate metrics.",
            "Interpret the graph and the result together.",
        ]),
        ("ARCHITECTURE", "The application is modular and packaged as a desktop executable.", [
            "CustomTkinter provides the interface.",
            "scikit-learn and statsmodels provide the algorithms.",
            "Matplotlib displays embedded visualizations.",
            "PyInstaller generates the final .exe.",
        ]),
        ("SCREENSHOT", "The regression tab demonstrates a complete ML experiment.", [
            "The user controls input ranges and noise.",
            "The model estimates coefficients and intercept.",
            "R2 and RMSE explain model quality.",
        ]),
        ("SCREENSHOT", "The clustering tab shows unsupervised learning in 3D.", [
            "K-Means discovers groups without labels.",
            "The 3D view shows separation between clusters.",
            "The silhouette score summarizes cluster quality.",
        ]),
        ("SCREENSHOTS", "Classification and validation show performance comparison.", [
            "Random Forest explains accuracy and feature importance.",
            "Cross-validation reduces dependence on one split.",
            "The application compares several models in one place.",
        ]),
        ("MODULES", "The six tabs cover the required machine learning families.", [
            "Regression: continuous prediction.",
            "Clustering 3D: unsupervised grouping.",
            "Random Forest: ensemble classification.",
            "ARIMA, neural network and cross-validation complete the TP.",
        ]),
        ("METHODS", "The report also studies additional algorithms.", [
            "SVM, Naive Bayes, Gradient Boosting and PCA.",
            "DBSCAN, Apriori, Hidden Markov Models and Autoencoders.",
            "Each method is explained with use cases, strengths and limits.",
        ]),
        ("APPLICATIONS", "The same ML logic applies to three options.", [
            "Development: bug prediction, logs, DevOps monitoring.",
            "Cybersecurity: anomaly detection, phishing and threat classification.",
            "Artificial Intelligence: prediction, segmentation and decision support.",
        ]),
        ("LIMITS", "Machine learning needs evaluation and human interpretation.", [
            "Random data is useful for learning but less realistic than production data.",
            "Overfitting, drift and poor data quality can reduce reliability.",
            "Metrics must be interpreted with graphs and domain knowledge.",
        ]),
        ("7-MIN PLAN", "The oral presentation is timed to stay under seven minutes.", [
            "1 minute: context, objective and AI/ML foundation.",
            "3 minutes: application workflow, architecture and screenshots.",
            "2 minutes: methods, domains, limits and perspectives.",
            "1 minute: conclusion and questions.",
        ]),
        ("CONCLUSION", "The final project is a complete academic and professional package.", [
            "The report proves scientific understanding.",
            "The executable proves implementation.",
            "The presentation and video prove communication.",
            "GitHub or LinkedIn publication turns it into a portfolio asset.",
        ]),
    ]

    for idx, (kicker, title, bullets) in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PptRGB(248, 250, 252)

        if idx == 1:
            hero = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(0), PptInches(10), PptInches(7.5))
            hero.fill.solid()
            hero.fill.fore_color.rgb = PptRGB(15, 23, 42)
            hero.line.fill.background()
            strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(0), PptInches(0.18), PptInches(7.5))
            strip.fill.solid()
            strip.fill.fore_color.rgb = PptRGB(56, 189, 248)
            strip.line.fill.background()
            title_box = slide.shapes.add_textbox(PptInches(0.75), PptInches(0.95), PptInches(8.2), PptInches(2.0))
            title_box.text_frame.word_wrap = True
            p = title_box.text_frame.paragraphs[0]
            p.text = title
            p.font.size = PptPt(28)
            p.font.bold = True
            p.font.color.rgb = PptRGB(255, 255, 255)
            add_bullet_box(slide, bullets, 0.78, 3.45, 5.45, 1.7, 14)
            for shape in slide.shapes:
                if shape.has_text_frame and shape.left < PptInches(6.2):
                    for para in shape.text_frame.paragraphs:
                        para.font.color.rgb = PptRGB(226, 232, 240)
            add_metric(slide, "6", "ML modules", 6.75, 3.2, PptRGB(37, 99, 235))
            add_metric(slide, ".exe", "desktop delivery", 6.75, 4.35, PptRGB(22, 163, 74))
            add_footer(slide, idx)
            continue

        add_slide_title(slide, kicker, title)

        if kicker == "TABLE OF CONTENT":
            colors = [PptRGB(37, 99, 235), PptRGB(22, 163, 74), PptRGB(249, 115, 22), PptRGB(124, 58, 237), PptRGB(14, 165, 233), PptRGB(220, 38, 38)]
            coords = [(0.75, 1.65), (3.45, 1.65), (6.15, 1.65), (0.75, 3.75), (3.45, 3.75), (6.15, 3.75)]
            for (x, y), item, color in zip(coords, bullets, colors):
                number, body = item.split(".", 1)
                if ":" in body:
                    heading, detail = body.strip().split(":", 1)
                else:
                    heading, detail = body.strip(), "Close the project story and prepare for questions."
                add_card(slide, f"{number}. {heading.strip()}", detail.strip(), x, y, 2.35, 1.55, color)
        elif kicker == "WORKFLOW":
            add_bullet_box(slide, bullets, 0.75, 1.75, 4.4, 2.5)
            slide.shapes.add_picture(str(ASSETS / "ml_workflow.png"), PptInches(5.25), PptInches(2.0), width=PptInches(4.1))
        elif kicker == "ARCHITECTURE":
            add_bullet_box(slide, bullets, 0.75, 1.75, 4.4, 2.5)
            slide.shapes.add_picture(str(ASSETS / "app_architecture.png"), PptInches(5.25), PptInches(1.95), width=PptInches(4.15))
        elif title == "The regression tab demonstrates a complete ML experiment.":
            add_bullet_box(slide, bullets, 0.75, 1.75, 3.2, 3.4, 12)
            add_image_frame(slide, screenshots / "screenshot_regression.png", 4.2, 1.85, 5.1, 3.1)
        elif title == "The clustering tab shows unsupervised learning in 3D.":
            add_bullet_box(slide, bullets, 0.75, 1.75, 3.2, 3.4, 12)
            add_image_frame(slide, screenshots / "screenshot_clustering.png", 4.2, 1.85, 5.1, 3.1)
        elif kicker == "SCREENSHOTS":
            add_image_frame(slide, screenshots / "screenshot_random_forest.png", 0.8, 1.85, 4.1, 2.5)
            add_image_frame(slide, screenshots / "screenshot_validation.png", 5.15, 1.85, 4.1, 2.5)
            add_bullet_box(slide, bullets, 1.0, 4.95, 8.1, 1.15, 11)
        elif kicker in ("MODULES", "METHODS"):
            labels = [b.split(":")[0] for b in bullets]
            bodies = [b.split(":", 1)[1].strip() if ":" in b else b for b in bullets]
            coords = [(0.75, 1.85), (3.25, 1.85), (0.75, 4.05), (3.25, 4.05)]
            colors = [PptRGB(37, 99, 235), PptRGB(22, 163, 74), PptRGB(249, 115, 22), PptRGB(124, 58, 237)]
            for (x, y), label, body, color in zip(coords, labels, bodies, colors):
                add_card(slide, label, body, x, y, 2.25, 1.55, color)
            add_side_panel(slide, "Defense focus", ["Explain what it does", "Read the metric", "Connect to the report"], PptRGB(37, 99, 235))
        elif kicker == "7-MIN PLAN":
            add_bullet_box(slide, bullets, 0.75, 1.75, 5.1, 3.8, 13)
            add_side_panel(slide, "Target", ["15 slides", "About 28 sec each", "Max 7 minutes"], PptRGB(22, 101, 52))
        else:
            add_bullet_box(slide, bullets, 0.75, 1.75, 5.25, 3.6, 13)
            panel_titles = {
                "REQUIREMENTS": ("Validation", ["Executable", "Report", "Presentation"]),
                "FOUNDATION": ("Core idea", ["Data", "Model", "Decision"]),
                "APPLICATIONS": ("Three options", ["Development", "Cybersecurity", "AI"]),
                "LIMITS": ("Risks", ["Data", "Overfitting", "Drift"]),
                "CONCLUSION": ("Final message", ["Research", "Implementation", "Communication"]),
            }
            heading, lines = panel_titles.get(kicker, ("Key point", ["Theory", "Practice", "Evaluation"]))
            add_side_panel(slide, heading, lines, PptRGB(15, 23, 42))

        add_footer(slide, idx)

    prs.save(DECK)


if __name__ == "__main__":
    create_assets()
    build_report()
    build_deck()
    print(REPORT)
    print(DECK)
